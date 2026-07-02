from __future__ import annotations

import io
import uuid
from typing import Any

import aiofiles
from fastapi import UploadFile
from sqlalchemy.ext.asyncio import AsyncSession

from app.ai.providers import get_ai_provider
from app.core.config import settings
from app.core.logging import get_logger
from app.models.knowledge_document import DocumentChunk, DocumentStatus, KnowledgeDocument

logger = get_logger(__name__)

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks


async def extract_text_from_file(file: UploadFile) -> str:
    content = await file.read()
    filename = file.filename or ""

    if filename.endswith(".txt"):
        return content.decode("utf-8", errors="ignore")

    if filename.endswith(".pdf"):
        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(io.BytesIO(content))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception as e:
            logger.error("PDF extraction error", error=str(e))
            return ""

    if filename.endswith(".docx"):
        try:
            import docx
            doc = docx.Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            logger.error("DOCX extraction error", error=str(e))
            return ""

    if filename.endswith(".csv"):
        return content.decode("utf-8", errors="ignore")

    if filename.endswith((".xlsx", ".xls")):
        try:
            import pandas as pd
            df = pd.read_excel(io.BytesIO(content))
            return df.to_string()
        except Exception as e:
            logger.error("Excel extraction error", error=str(e))
            return ""

    if filename.endswith((".pptx",)):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
        except Exception as e:
            logger.error("PPTX extraction error", error=str(e))
            return ""

    return content.decode("utf-8", errors="ignore")


class KnowledgeBaseProcessor:
    def __init__(self, session: AsyncSession) -> None:
        self.session = session
        self.provider = get_ai_provider()

    async def process_upload(self, file: UploadFile, uploaded_by: uuid.UUID | None = None) -> KnowledgeDocument:
        filename = file.filename or "document"
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "txt"

        doc = KnowledgeDocument(
            title=filename,
            file_name=filename,
            doc_type=ext,
            status=DocumentStatus.PROCESSING,
            uploaded_by=uploaded_by,
        )
        self.session.add(doc)
        await self.session.flush()

        try:
            text = await extract_text_from_file(file)
            await self._create_chunks(doc, text)
            doc.status = DocumentStatus.READY
        except Exception as e:
            logger.error("Document processing failed", error=str(e), doc_id=str(doc.id))
            doc.status = DocumentStatus.FAILED
            doc.error_message = str(e)

        self.session.add(doc)
        await self.session.flush()
        return doc

    async def process_url(self, url: str, title: str = "", uploaded_by: uuid.UUID | None = None) -> KnowledgeDocument:
        doc = KnowledgeDocument(
            title=title or url,
            source_url=url,
            doc_type="url",
            status=DocumentStatus.PROCESSING,
            uploaded_by=uploaded_by,
        )
        self.session.add(doc)
        await self.session.flush()

        try:
            import httpx
            async with httpx.AsyncClient(timeout=30) as client:
                resp = await client.get(url, follow_redirects=True)
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(resp.text, "html.parser")
                for tag in soup(["script", "style", "nav", "footer"]):
                    tag.decompose()
                text = soup.get_text(separator="\n", strip=True)
            await self._create_chunks(doc, text)
            doc.status = DocumentStatus.READY
        except Exception as e:
            logger.error("URL processing failed", error=str(e))
            doc.status = DocumentStatus.FAILED
            doc.error_message = str(e)

        self.session.add(doc)
        await self.session.flush()
        return doc

    async def _create_chunks(self, doc: KnowledgeDocument, text: str) -> None:
        chunks = chunk_text(text)
        doc.total_chunks = len(chunks)

        for i, chunk_content in enumerate(chunks):
            embedding = await self.provider.embed_single(chunk_content)
            chunk = DocumentChunk(
                document_id=doc.id,
                chunk_index=i,
                content=chunk_content,
                embedding=embedding,
                token_count=len(chunk_content.split()),
            )
            self.session.add(chunk)

        await self.session.flush()
